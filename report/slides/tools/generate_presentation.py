from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
SLIDES_DIR = ROOT / "slides"
TEMPLATE = max(SLIDES_DIR.glob("*.pptx"), key=lambda path: path.stat().st_size)
OUTPUT = SLIDES_DIR / "光电混合运动想象脑机接口_复赛答辩.pptx"

NAVY = RGBColor(24, 47, 78)
BLUE = RGBColor(42, 101, 160)
RED = RGBColor(198, 64, 62)
GREEN = RGBColor(56, 132, 91)
INK = RGBColor(32, 43, 56)
MUTED = RGBColor(91, 104, 119)
LIGHT = RGBColor(242, 245, 248)
WHITE = RGBColor(255, 255, 255)


@dataclass(frozen=True)
class SlideSpec:
    title: str
    bullets: tuple[str, ...]
    images: tuple[str, ...]
    note: str
    accent: RGBColor = BLUE
    section: str = ""


def image(relative: str) -> str:
    return str(ROOT / Path(relative))


SPECS = [
    SlideSpec(
        "光电混合运动想象脑机接口",
        ("FBCSP", "经验库特化", "Photonic Candidate Scan", "参赛编号：现场填写"),
        (image("template/figure/competition/mi_bci_task_overview.png"),),
        "各位评委好，我们的作品是“光电混合运动想象脑机接口”。项目使用 FBCSP 提取可解释的运动想象脑电特征，通过小型 MLP 和经验库完成会话特化，再用低位宽 photonic candidate scan 扫描多个候选分类头。系统输出左手、右手、脚三类命令；结果不可靠时输出拒识。",
        RED,
        "作品概览",
    ),
    SlideSpec(
        "赛题要求与作品定位",
        ("复赛：创新 AI 应用", "光计算占比 ≥ 50%", "方向：MI-BCI", "证据：公开数据回放", "边界：软件量化模拟"),
        (image("template/figure/competition/mi_bci_task_overview.png"),),
        "复赛要求基于光计算平台开发创新 AI 应用，并要求光计算占比不低于 50%。我们选择运动想象脑机接口作为应用方向，将适合矩阵和向量运算的前向线性计算统一交给 photonic backend。当前结果来自公开 EEG 数据回放和软件量化模拟，不是 Cyton 真实在线闭环，也不是真实光芯片功耗实测。",
        RED,
        "作品概览",
    ),
    SlideSpec(
        "为什么选择运动想象 EEG",
        ("低信噪比", "小样本", "非平稳", "跨被试差异", "拒识：避免误触发"),
        (image("template/figure/competition/eeg_mi_signal_concepts.png"), image("template/figure/competition/eeg_1020_electrodes.jpg")),
        "EEG 信号幅值弱、信噪比低，数据规模又远小于图像任务，并且会受到眨眼、肌电、电极接触和状态变化影响。运动想象主要关注 8 到 30 赫兹附近的 mu 和 beta 节律，提示出现后常发生 ERD 功率下降。不同被试的有效频段和空间激活模式差异明显，所以系统既要保持小样本稳定性，也要允许会话特化，还要通过拒识减少错误控制命令。",
        GREEN,
        "问题定义",
    ),
    SlideSpec(
        "核心流程",
        ("EEG → FBCSP", "→ Small MLP Embedding", "→ Experience Library", "→ Photonic Scan", "→ Reject / Command"),
        (image("template/figure/competition/end_to_end_inference_chain.drawio.png"),),
        "系统首先从 EEG 窗口提取 FBCSP 时频和空间特征，再由小型 MLP 映射到紧凑 embedding。少量校准窗口用于查询经验库，选择并加权多个候选线性头；随后 photonic scan 计算候选得分，数字端完成 softmax、概率融合和拒识判断。这里 FBCSP 是主干，小型网络只改善特征空间，并不是纯深度学习端到端 BCI。",
        BLUE,
        "总体方案",
    ),
    SlideSpec(
        "训练 · 校准 · 在线推理",
        ("训练：FBCSP / MLP / 经验库", "校准：42 窗口 / Top-K", "推理：单窗口前向", "Train ≠ Calibration ≠ Evaluation"),
        (
            image("slides/assets/stage_offline_training.png"),
            image("slides/assets/stage_session_calibration.png"),
            image("slides/assets/stage_online_inference.png"),
        ),
        "离线训练阶段拟合 FBCSP、特征选择、标准化和小型 MLP，并构建经验库。校准阶段只使用少量目标会话窗口生成 embedding，查询 top-K 候选并冻结融合权重和拒识阈值。在线推理阶段每次只处理一个新窗口，不再训练模型。BCICIV 主线使用 42 个校准窗口，剩余 518 个窗口独立评估，避免训练、校准和评估之间的数据泄漏。",
        BLUE,
        "总体方案",
    ),
    SlideSpec(
        "FBCSP：可解释的线性主干",
        ("6 个运动相关频带", "OVR CSP 空间投影", "Log-variance", "72 → 32 维", "小样本 · 线性计算密集"),
        (image("template/figure/competition/fbcsp_matrix_flow.drawio.png"),),
        "FBCSP 将 EEG 分到六个运动相关频带，对每个类别执行 one-versus-rest CSP 空间投影，再计算归一化对数方差。原始特征共 72 维，通过 Fisher 分数筛选为 32 维。它适合小样本，频带和空间模式也便于解释。更重要的是，CSP 投影、协方差和 Gram 矩阵包含大量规则线性计算，因此 FBCSP 不只是 baseline，也是光计算接管的重要主干。",
        GREEN,
        "算法设计",
    ),
    SlideSpec(
        "个体差异驱动经验库特化",
        ("2 Anchor Heads", "64 Bootstrap Heads", "42 校准窗口", "Top-K = 8", "候选模型 ≠ 数据缓存"),
        (image("template/figure/competition/specialization_logic.png"), image("hybrid-photonic-mi-bci/artifacts/figures/fbcsp_design/experience_photonic/mainline_retrieval_weights.png")),
        "不同被试的节律频段、空间激活和想象策略并不一致，所以我们没有只依赖一个固定全局模型。当前经验库包含两个全局 anchor heads 和 64 个 bootstrap heads。42 个校准窗口产生当前会话 embedding，并综合距离、训练质量、校准准确率和置信度选择 top-8。图中的两个 anchor 权重较高，其他候选提供会话特化补充。经验库存储的是候选模型和质量统计，不是普通样本缓存。",
        GREEN,
        "算法设计",
    ),
    SlideSpec(
        "Photonic Candidate Scan",
        ("Head Bank：(8, 3, 33)", "Tile：2 × 8", "10 Tiles / Candidate", "80 Tiles / Window", "输出：Candidate Scores"),
        (image("template/figure/competition/photonic_tile_mapping.drawio.png"), image("hybrid-photonic-mi-bci/artifacts/figures/fbcsp_design/experience_photonic/photonic_tile_schedule.png")),
        "八个候选头堆叠后的 weights shape 是 8 乘 3 乘 33，其中 33 包括 32 维 embedding 和一个偏置常数。物理 tile 为 2 乘 8，因此每个候选需要 2 个行块和 5 个列块，也就是 10 次 tile；每个窗口共 80 次。Scan 输出的是八组候选分类分数，不是最终类别。数字端再执行 softmax、检索权重融合和拒识。这里的 scan 是候选分类头扫描，不是 Goertzel 或频率扫描。",
        RED,
        "光计算映射",
    ),
    SlideSpec(
        "统一后端接口",
        ("MatrixOps：matmul / einsum", "SignalOps：CAR / SOS", "TiledMVM：Candidate Bank", "上层算法 ↔ 执行核解耦", "当前：软件模拟"),
        (image("template/figure/competition/system_block_diagram.png"),),
        "MatrixOpsBackend 负责通用 matmul 和 einsum，覆盖 CSP、标准化、MLP、LDA、距离项和概率融合。SignalOpsBackend 负责带通道轴或时间轴语义的 CAR 和 SOS 滤波。TiledMVMBackend 负责候选 head bank 的分块扫描和 tile 计数。上层算法只依赖接口，不直接依赖硬件。当前执行核仍是软件模拟，未来可以在不改动决策逻辑的情况下替换为官方模拟器或真实芯片驱动。",
        BLUE,
        "工程接口",
    ),
    SlideSpec(
        "量化策略：逻辑精度 ≠ 物理位宽",
        ("Candidate：单次 4-bit", "Input：uint4 [0,15]", "Weight：int4 [-8,7]", "其他路径：8-bit Logic", "Radix-16 / 4-bit Slices"),
        (image("template/figure/competition/tile_schedule.png"),),
        "候选扫描直接使用单次 4-bit：输入是 0 到 15 的 uint4，权重是负 8 到 7 的 int4。其他前向线性路径默认保留 8-bit 逻辑精度，但会采用 radix-16 拆成多个 4-bit slice，再按位权重构。因此 8-bit 逻辑精度不是一次原生 8-bit 光计算调用。当前采用运行时动态量化，不是 QAT。后续需要逐算子搜索最低可用精度，联合比较准确率、拒识率、slice 数、噪声、延迟和能耗。",
        RED,
        "光计算映射",
    ),
    SlideSpec(
        "实验协议",
        ("BCICIV：3 类 MI", "Train：840", "Calibration：42", "Evaluation：518", "BNCI：9 Subjects / 1296"),
        (image("hybrid-photonic-mi-bci/artifacts/figures/bnci2014_004_personalization/bnci_subject_line_comparison.png"),),
        "BCICIV 用于三类运动想象主流程验证。训练集合用于 FBCSP、MLP 和经验库构建，42 个窗口只用于校准，518 个窗口只用于最终回放评估。BNCI 数据包含 9 位被试，用 session 1 和 2 作为历史数据，session 3 作为目标会话，共保留 1296 个评估窗口。Reject 不是数据集原生类别，而是系统根据置信度额外产生的安全状态。",
        BLUE,
        "实验验证",
    ),
    SlideSpec(
        "BCICIV 主线结果",
        ("Accuracy  70.46%", "Balanced  73.24%", "Accepted  72.42%", "Reject  2.70%", "Evaluation  518"),
        (image("hybrid-photonic-mi-bci/artifacts/figures/fbcsp_design/experience_photonic/mainline_confusion.png"), image("hybrid-photonic-mi-bci/artifacts/figures/fbcsp_design/experience_photonic/mainline_cumulative_metrics.png")),
        "在 BCICIV 的 518 个独立评估窗口上，主线 command accuracy 为 70.46%，balanced accuracy 为 73.24%，accepted accuracy 为 72.42%，reject rate 为 2.70%。混淆矩阵显示 left 和 right 仍有明显混淆，foot 样本相对较少。累计曲线用于观察回放稳定性。拒识率不能单独追求更高，它需要与 accepted accuracy 和可输出命令数量共同分析。",
        GREEN,
        "实验验证",
    ),
    SlideSpec(
        "三线对比与光计算账本",
        ("LDA  71.25%", "Small MLP  72.32%", "Mainline  70.46%", "Forward Share  100%", "Inference MAC  342.105 M"),
        (image("hybrid-photonic-mi-bci/artifacts/figures/fbcsp_design/summary/design_line_summary.png"), image("hybrid-photonic-mi-bci/artifacts/figures/fbcsp_design/compute_accounting/compute_accounting_summary.png")),
        "同一协议下，LDA baseline 的准确率是 71.25%，小型 MLP 是 72.32%，主线是 70.46%。因此我们不声称经验库已经稳定提升平均准确率；当前已验证的是校准、特化、多候选扫描和拒识的完整路径。按 forward-only 账本，342.105 M 前向线性 MAC 全部路由到 photonic backend，接管比例为 100%。这个数字是接口和 MAC-equivalent 口径，不等于真实芯片功耗占比或实际加速比。",
        RED,
        "实验验证",
    ),
    SlideSpec(
        "创新 · 边界 · 下一步",
        ("FBCSP 光计算主干", "经验库会话特化", "多候选低位宽扫描", "未完成：Cyton 在线闭环", "下一步：真实芯片 / 混合精度"),
        (image("template/figure/competition/specialization_logic.png"), image("template/figure/competition/photonic_tile_mapping.drawio.png")),
        "项目有三个主要设计点：第一，FBCSP 既提供可解释的小样本主干，也暴露大量可接管线性计算；第二，经验库针对 MI-EEG 个体差异完成少量校准后的会话特化；第三，利用多候选结构设计低位宽 photonic scan，而不是只替换一个矩阵乘。目前 Cyton 在线闭环和真实光芯片实测尚未完成。下一步将接入真实驱动，开展逐算子混合精度、噪声、延迟和功耗测试，并完成真实被试在线验证。谢谢各位评委。",
        NAVY,
        "总结",
    ),
    SlideSpec(
        "备用｜小型 MLP 与 Embedding",
        ("32 → 64 → 32", "Training Loss", "Embedding PCA", "Replay Trace", "Confusion"),
        (image("hybrid-photonic-mi-bci/artifacts/figures/fbcsp_design/small_network/small_network_training_embedding.png"),),
        "这页用于回答小型网络是否正常训练。左上显示 loss 下降和训练准确率上升；右上是 32 维 embedding 的 PCA 二维投影，仅用于可视化，不是实际分类空间；下方是独立回放的滚动指标和混淆矩阵。训练准确率较高，因此必须结合回放结果判断泛化。这里使用当前生成的 small-network 图，而不是旧版 embedding diagnostics。",
        MUTED,
        "备用答疑",
    ),
    SlideSpec(
        "备用｜BNCI 个体结果",
        ("9 Subjects", "Mainline  74.92%", "Accepted  76.22%", "Reject  1.62%", "Subject-wise 差异"),
        (image("hybrid-photonic-mi-bci/artifacts/figures/bnci2014_004_personalization/bnci_subject_line_comparison.png"), image("hybrid-photonic-mi-bci/artifacts/figures/bnci2014_004_personalization/bnci_design_line_summary.png")),
        "BNCI 用来观察跨被试和目标会话差异。九位被试的主线平均准确率为 74.92%，accepted accuracy 为 76.22%，reject rate 为 1.62%。LDA 平均值仍略高，因此经验库价值需要看 subject-wise 收益，而不能只看总体平均。后续会对校准窗口数、top-K 和经验库规模做消融。",
        MUTED,
        "备用答疑",
    ),
    SlideSpec(
        "备用｜工程与复现",
        ("30 Tests Passed", "JSON / NPZ", "Single-window Inference", "Pure Runtime", "Reproducible Figures"),
        (image("hybrid-photonic-mi-bci/artifacts/figures/fbcsp_design/system/system_block_diagram_detailed.png"),),
        "工程提供完整三线、单主线、单窗口和 BNCI 四类运行入口。JSON 保存 summary 和计算账本，NPZ 保存概率、embedding、rolling、cumulative 和混淆矩阵原始数组，绘图脚本可以重新生成所有结果图。当前回归测试为 30 项通过。Pure runtime 将部署前向与训练、评估和绘图分离，便于后续替换真实执行后端。",
        MUTED,
        "备用答疑",
    ),
]


def remove_all_slides(prs: Presentation) -> None:
    for slide_id in list(prs.slides._sldIdLst):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def add_text(slide, text, x, y, w, h, *, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_picture_contain(slide, path: str, x, y, w, h):
    from PIL import Image

    file_path = Path(path)
    if not file_path.exists():
        raise FileNotFoundError(file_path)
    with Image.open(file_path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    px, py = x + (w - pw) // 2, y + (h - ph) // 2
    return slide.shapes.add_picture(str(file_path), px, py, width=pw, height=ph)


def add_frame(slide, x, y, w, h, accent):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(216, 223, 231)
    shape.line.width = Pt(1)
    return shape


def add_header(slide, spec: SlideSpec, number: int, total: int):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.16))
    band.fill.solid()
    band.fill.fore_color.rgb = spec.accent
    band.line.fill.background()
    add_text(slide, spec.section.upper(), Inches(0.62), Inches(0.27), Inches(2.9), Inches(0.3), size=10, color=spec.accent, bold=True)
    add_text(slide, spec.title, Inches(0.62), Inches(0.58), Inches(11.6), Inches(0.62), size=27, color=NAVY, bold=True)
    add_text(slide, f"{number:02d} / {total:02d}", Inches(11.7), Inches(0.3), Inches(0.95), Inches(0.28), size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.24), Inches(12.05), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(222, 228, 234)
    line.line.fill.background()


def add_bullets(slide, bullets, x, y, w, h, accent, *, size=19):
    frame_shape = slide.shapes.add_textbox(x, y, w, h)
    frame = frame_shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.03)
    for idx, item in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(10)
        p.line_spacing = 1.05
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.bold = idx == 0 and len(bullets) <= 4
        p.font.color.rgb = INK
        p._p.get_or_add_pPr().insert(0, p._p._new_buChar()) if False else None
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - Inches(0.14), y + Inches(0.09) + idx * Inches(0.55), Inches(0.05), Inches(0.22))
        marker.fill.solid()
        marker.fill.fore_color.rgb = accent
        marker.line.fill.background()
    return frame_shape


def add_notes(slide, note: str):
    try:
        text_frame = slide.notes_slide.notes_text_frame
        text_frame.text = note
    except Exception:
        pass


def build_slide(prs: Presentation, spec: SlideSpec, number: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE
    add_header(slide, spec, number, total)

    left_x, left_y, left_w, left_h = Inches(0.72), Inches(1.55), Inches(3.15), Inches(5.22)
    image_x, image_y, image_w, image_h = Inches(4.08), Inches(1.47), Inches(8.55), Inches(5.45)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, left_y, left_w, left_h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT
    panel.line.fill.background()
    add_bullets(slide, spec.bullets, left_x + Inches(0.28), left_y + Inches(0.35), left_w - Inches(0.45), left_h - Inches(0.6), spec.accent, size=17 if len(spec.bullets) >= 5 else 19)

    if len(spec.images) == 1:
        add_frame(slide, image_x, image_y, image_w, image_h, spec.accent)
        add_picture_contain(slide, spec.images[0], image_x + Inches(0.10), image_y + Inches(0.10), image_w - Inches(0.20), image_h - Inches(0.20))
    elif len(spec.images) == 2:
        gap = Inches(0.16)
        half_h = (image_h - gap) // 2
        for idx, pic in enumerate(spec.images[:2]):
            py = image_y + idx * (half_h + gap)
            add_frame(slide, image_x, py, image_w, half_h, spec.accent)
            add_picture_contain(slide, pic, image_x + Inches(0.08), py + Inches(0.08), image_w - Inches(0.16), half_h - Inches(0.16))
    else:
        gap = Inches(0.10)
        third_h = (image_h - gap * 2) // 3
        for idx, pic in enumerate(spec.images[:3]):
            py = image_y + idx * (third_h + gap)
            add_frame(slide, image_x, py, image_w, third_h, spec.accent)
            add_picture_contain(slide, pic, image_x + Inches(0.06), py + Inches(0.04), image_w - Inches(0.12), third_h - Inches(0.08))

    footer = add_text(slide, "第十届全国大学生集成电路创新创业大赛｜匿名答辩材料", Inches(0.7), Inches(7.08), Inches(7.4), Inches(0.22), size=8, color=MUTED)
    add_notes(slide, spec.note)
    return slide


def main() -> None:
    shutil.copy2(TEMPLATE, OUTPUT)
    prs = Presentation(str(OUTPUT))
    remove_all_slides(prs)
    total = len(SPECS)
    for idx, spec in enumerate(SPECS, start=1):
        build_slide(prs, spec, idx, total)
    prs.core_properties.title = "光电混合运动想象脑机接口复赛答辩"
    prs.core_properties.subject = "第十届全国大学生集成电路创新创业大赛"
    prs.core_properties.author = "参赛团队"
    prs.core_properties.keywords = "MI-BCI,FBCSP,photonic scan"
    prs.core_properties.comments = "匿名答辩材料"
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
