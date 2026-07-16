"""Generate the replay decision-flow asset and insert it after Candidate Scan."""

from copy import deepcopy
from pathlib import Path
from xml.etree.ElementTree import Element, ElementTree, SubElement

import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch, Polygon
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[2]
ASSETS = ROOT / "slides" / "assets"
PNG = ASSETS / "replay_decision_flow.png"
DRAWIO = ASSETS / "replay_decision_flow.drawio"
SOURCE = ROOT / "slides" / "ppt_v1.2.pptx"
OUTPUT = ROOT / "slides" / "ppt_v1.2_含replay决策页.pptx"

plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False

INK = "#203047"
MUTED = "#667487"
BLUE = "#3478B8"
GREEN = "#3C9064"
ORANGE = "#E58A2B"
PURPLE = "#7656A8"
RED = "#C95757"
TEAL = "#2F8791"
LINE = "#CAD5E1"
LIGHT = "#F5F8FB"


def rounded(ax, x, y, w, h, fill, edge=LINE, lw=1.35, radius=.07):
    ax.add_patch(FancyBboxPatch(
        (x, y), w, h, boxstyle=f"round,pad=0.02,rounding_size={radius}",
        facecolor=fill, edgecolor=edge, linewidth=lw,
    ))


def arrow(ax, start, end, color=MUTED, *, dashed=False, rad=0):
    connection = f"arc3,rad={rad}" if rad else "arc3,rad=0"
    ax.add_patch(FancyArrowPatch(
        start, end, connectionstyle=connection, arrowstyle="-|>", mutation_scale=13,
        linewidth=1.55, linestyle="--" if dashed else "-", color=color,
    ))


def flow_box(ax, x, y, w, h, color, title, detail, *, light=False, label=None):
    rounded(ax, x, y, w, h, LIGHT if light else color, color, 1.35)
    if label:
        ax.text(x + .15, y + h - .17, label, ha="left", va="center", fontsize=7.5,
                color=color if light else "white", weight="bold")
    ax.text(x + w / 2, y + h * .60, title, ha="center", va="center", fontsize=11.8,
            color=INK if light else "white", weight="bold")
    ax.text(x + w / 2, y + h * .28, detail, ha="center", va="center", fontsize=8.4,
            color=MUTED if light else "white", linespacing=1.25)


def generate_png():
    fig, ax = plt.subplots(figsize=(13.33, 7.5))
    ax.set(xlim=(0, 13.33), ylim=(0, 7.5))
    ax.axis("off")
    fig.patch.set_facecolor("white")

    ax.text(.48, 7.05, "Replay 决策：从候选扫描到安全控制输出",
            fontsize=20, weight="bold", color=INK)
    ax.text(.49, 6.63,
            "候选评分由 photonic scan 完成；softmax、概率融合与拒识判断保留在数字端",
            fontsize=10.5, color=MUTED)

    # Main handoff: scan score -> digital decision path.
    rounded(ax, .47, 4.48, 12.39, 1.66, "#F8FAFC", "#D9E2EC", 1.0)
    ax.text(.70, 5.84, "每个 replay 窗口的候选决策链", fontsize=10.7,
            weight="bold", color=INK)
    ax.text(.70, 5.54, "候选扫描输出的是 K 组 3 类得分，不直接等同于最终控制命令",
            fontsize=8.6, color=MUTED)

    items = [
        (.70, 1.80, PURPLE, "Photonic scan", "Top-K 候选头\n并行评分", False, "光计算端"),
        (2.83, 1.72, ORANGE, "Candidate scores", "s₁ … sᴷ\n每个候选 3 类得分", False, "读回结果"),
        (4.86, 1.70, BLUE, "+ bias", "数字端补偿 bᵢ\n形成校正 logits", False, "数字端"),
        (6.88, 1.58, TEAL, "Softmax", "pᵢ = softmax(sᵢ + bᵢ)\n每候选归一化", False, "数字端"),
        (8.78, 2.10, GREEN, "概率融合", "p = Σ αᵢ · pᵢ\n按检索权重 αᵢ 融合", False, "数字端"),
    ]
    for x, w, color, title, detail, light, label in items:
        flow_box(ax, x, 4.79, w, .88, color, title, detail, light=light, label=label)
    for start, end, color in [
        ((2.53, 5.23), (2.80, 5.23), PURPLE),
        ((4.58, 5.23), (4.83, 5.23), ORANGE),
        ((6.58, 5.23), (6.85, 5.23), BLUE),
        ((8.49, 5.23), (8.75, 5.23), TEAL),
    ]:
        arrow(ax, start, end, color)

    # The three independent safety tests visibly meet at one decision point.
    ax.text(.70, 4.03, "Safety gate：三项均满足才允许产生控制指令", fontsize=10.7,
            weight="bold", color=RED)
    ax.text(.70, 3.73, "任一条件不满足，输出 reject；reject 是安全状态，不是第 4 类运动想象", fontsize=8.6,
            color=MUTED)
    gate_boxes = [
        (.70, BLUE, "Confidence", "max(p) ≥ τₚ\n最大类别概率", "概率足够集中？"),
        (3.27, ORANGE, "Margin", "p⁽¹⁾ − p⁽²⁾ ≥ τₘ\n第一、二类差距", "类别足够分离？"),
        (5.84, PURPLE, "Quality", "q < τq 且无伪迹\n坏道 / 饱和 / 肌电", "信号质量合格？"),
    ]
    for x, color, title, detail, prompt in gate_boxes:
        flow_box(ax, x, 2.32, 2.18, 1.00, color, title, detail, label=prompt)
    arrow(ax, (9.82, 4.78), (9.82, 3.41), GREEN)
    for x, color in [(.70, BLUE), (3.27, ORANGE), (5.84, PURPLE)]:
        arrow(ax, (x + 1.09, 2.30), (8.30, 1.84), color, rad=-.10)

    # Gate diamond and the two semantically distinct outputs.
    diamond = Polygon([(9.25, 2.36), (10.15, 1.84), (9.25, 1.32), (8.35, 1.84)],
                      closed=True, facecolor="#FFF4F4", edgecolor=RED, linewidth=1.6)
    ax.add_patch(diamond)
    ax.text(9.25, 1.94, "全部\n通过？", ha="center", va="center", fontsize=10.5,
            weight="bold", color=RED, linespacing=1.05)
    arrow(ax, (10.18, 1.84), (10.68, 1.84), GREEN)
    arrow(ax, (9.25, 1.28), (9.25, .81), RED)
    flow_box(ax, 10.72, 1.32, 2.04, 1.02, GREEN, "Command", "arg max p\n左手 / 右手 / 脚", label="是：生成控制指令")
    flow_box(ax, 7.92, .16, 2.66, .62, RED, "Reject", "不把低置信窗口转换为控制指令", label="否：安全拒识")
    ax.text(6.69, .47, "三类输出 + 拒识\n≠ 强制三分类", ha="right", va="center", fontsize=9.0,
            color=RED, weight="bold")

    rounded(ax, .47, .13, 6.13, .54, "#EEF3F8", LINE, 1.0)
    ax.text(3.535, .40, "融合后才评估 confidence / margin；质量异常可直接触发 reject",
            ha="center", va="center", fontsize=9.3, color=INK, weight="bold")

    fig.savefig(PNG, dpi=220, bbox_inches="tight", facecolor="white", pad_inches=.08)
    plt.close(fig)


BASE = "rounded=1;whiteSpace=wrap;html=1;arcSize=10;align=center;verticalAlign=middle;fontFamily=Microsoft YaHei;"


def vertex(root, ident, value, x, y, w, h, fill, stroke, font=INK, size=14, sw=1):
    style = BASE + f"fillColor={fill};strokeColor={stroke};strokeWidth={sw};fontColor={font};fontSize={size};"
    cell = SubElement(root, "mxCell", id=ident, value=value, style=style, vertex="1", parent="1")
    SubElement(cell, "mxGeometry", x=str(x), y=str(y), width=str(w), height=str(h), **{"as": "geometry"})


def edge(root, ident, source, target, color=MUTED, label="", dashed=False):
    style = ("edgeStyle=orthogonalEdgeStyle;rounded=0;html=1;strokeWidth=2;"
             f"strokeColor={color};endArrow=block;endFill=1;"
             f"{'dashed=1;' if dashed else ''}fontFamily=Microsoft YaHei;fontSize=11;fontColor={color};")
    cell = SubElement(root, "mxCell", id=ident, value=label, style=style,
                      edge="1", parent="1", source=source, target=target)
    SubElement(cell, "mxGeometry", relative="1", **{"as": "geometry"})


def generate_drawio():
    mx = Element("mxfile", host="app.diagrams.net", modified="2026-07-16T00:00:00.000Z", agent="Codex", version="24.7.17")
    diagram = SubElement(mx, "diagram", id="replay-decision", name="Page-1")
    model = SubElement(diagram, "mxGraphModel", dx="1365", dy="767", grid="1", gridSize="10", guides="1", tooltips="1", connect="1", arrows="1", fold="1", page="1", pageScale="1", pageWidth="1333", pageHeight="750", math="0", shadow="0")
    root = SubElement(model, "root")
    SubElement(root, "mxCell", id="0")
    SubElement(root, "mxCell", id="1", parent="0")
    vertex(root, "title", "<b>Replay 决策：从候选扫描到安全控制输出</b><br><font color='#667487' style='font-size:12px'>候选评分由 photonic scan 完成；softmax、概率融合与拒识判断保留在数字端</font>", 45, 20, 1240, 65, "none", "none", INK, 24, 0)
    vertex(root, "lane", "", 45, 105, 1240, 170, "#F8FAFC", "#D9E2EC", INK, 12, 1)
    vertex(root, "scan", "<b>Photonic scan</b><br><font style='font-size:11px'>Top-K 候选头并行评分</font>", 65, 165, 180, 85, PURPLE, PURPLE, "#FFFFFF")
    vertex(root, "score", "<b>Candidate scores</b><br><font style='font-size:11px'>s₁ … sᴷ；每个候选 3 类得分</font>", 285, 165, 180, 85, ORANGE, ORANGE, "#FFFFFF")
    vertex(root, "bias", "<b>+ bias</b><br><font style='font-size:11px'>数字端补偿 bᵢ，形成校正 logits</font>", 505, 165, 180, 85, BLUE, BLUE, "#FFFFFF")
    vertex(root, "softmax", "<b>Softmax</b><br><font style='font-size:11px'>pᵢ = softmax(sᵢ + bᵢ)</font>", 725, 165, 180, 85, TEAL, TEAL, "#FFFFFF")
    vertex(root, "fusion", "<b>概率融合</b><br><font style='font-size:11px'>p = Σ αᵢ · pᵢ</font>", 945, 165, 220, 85, GREEN, GREEN, "#FFFFFF")
    vertex(root, "gate-title", "<b>Safety gate：三项均满足才允许产生控制指令</b><br><font color='#667487' style='font-size:11px'>任一条件不满足，输出 reject；reject 是安全状态，不是第 4 类运动想象</font>", 65, 305, 700, 50, "none", "none", RED, 15, 0)
    vertex(root, "confidence", "<b>Confidence</b><br><font style='font-size:11px'>max(p) ≥ τₚ<br>最大类别概率</font>", 65, 390, 210, 105, BLUE, BLUE, "#FFFFFF")
    vertex(root, "margin", "<b>Margin</b><br><font style='font-size:11px'>p⁽¹⁾ − p⁽²⁾ ≥ τₘ<br>第一、二类差距</font>", 315, 390, 210, 105, ORANGE, ORANGE, "#FFFFFF")
    vertex(root, "quality", "<b>Quality</b><br><font style='font-size:11px'>q &lt; τq 且无伪迹<br>坏道 / 饱和 / 肌电</font>", 565, 390, 210, 105, PURPLE, PURPLE, "#FFFFFF")
    vertex(root, "gate", "<b>全部<br>通过？</b>", 850, 390, 135, 105, "#FFF4F4", RED, RED, 16, 2)
    vertex(root, "command", "<b>Command</b><br><font style='font-size:11px'>arg max p<br>左手 / 右手 / 脚</font>", 1050, 390, 200, 105, GREEN, GREEN, "#FFFFFF")
    vertex(root, "reject", "<b>Reject</b>　不把低置信窗口转换为控制指令", 720, 565, 380, 62, RED, RED, "#FFFFFF", 14, 1)
    vertex(root, "note", "<b>融合后才评估 confidence / margin；质量异常可直接触发 reject</b>", 45, 660, 620, 40, "#EEF3F8", LINE, INK, 13, 1)
    for ident, source, target, color in [
        ("a1", "scan", "score", PURPLE), ("a2", "score", "bias", ORANGE),
        ("a3", "bias", "softmax", BLUE), ("a4", "softmax", "fusion", TEAL),
        ("b1", "fusion", "confidence", GREEN), ("b2", "confidence", "gate", BLUE),
        ("b3", "margin", "gate", ORANGE), ("b4", "quality", "gate", PURPLE),
        ("b5", "gate", "command", GREEN), ("b6", "gate", "reject", RED),
    ]:
        edge(root, ident, source, target, color)
    ElementTree(mx).write(DRAWIO, encoding="utf-8", xml_declaration=True)


def add_to_deck():
    prs = Presentation(SOURCE)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = __import__("pptx").dml.color.RGBColor(255, 255, 255)
    slide.shapes.add_picture(str(PNG), 0, 0, width=prs.slide_width, height=prs.slide_height)

    # New slide is initially last; move it directly after the existing Candidate Scan page (slide 10).
    slide_id_list = prs.slides._sldIdLst
    new_id = slide_id_list[-1]
    slide_id_list.remove(new_id)
    slide_id_list.insert(10, new_id)
    prs.save(OUTPUT)


if __name__ == "__main__":
    ASSETS.mkdir(parents=True, exist_ok=True)
    generate_png()
    generate_drawio()
    add_to_deck()
